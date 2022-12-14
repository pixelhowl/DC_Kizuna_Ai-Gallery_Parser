"""PPT automation"""
import os
import re

# pylint: disable=protected-access, unused-variable, logging-fstring-interpolation
import pandas as pd
from pptx import Presentation
# from pptx.chart.data import ChartData
from pptx.dml.color import RGBColor
from pptx.util import Cm

from utils import database, logging, paths, strings

TEMPLATE = Presentation(paths.POWERPOINT_FILE)
TEMPLATE_MONTH = 7


def replace_text_only(paragraph, target_text, color):
    # cur_text = paragrah.runs[0].text
    # new_text = cur_text.replace(str(search_str), str(repl_str))
    for idx, run in enumerate(paragraph.runs):
        if idx > 0:
            p = paragraph._p
            p.remove(run._r)
    paragraph.runs[0].font.color.rgb = RGBColor(*color)
    paragraph.runs[0].text = target_text


def replace_picture(slide, shape, image_path):
    with open(image_path, "rb") as file_obj:
        r_img_blob = file_obj.read()
    img_pic = shape._pic
    img_rid = img_pic.xpath("./p:blipFill/a:blip/@r:embed")[0]
    img_part = slide.part.related_part(img_rid)
    img_part._blob = r_img_blob


def replace_texts_with_diff_color(paragraph, target_texts, colors):
    for idx, run in enumerate(paragraph.runs):
        if idx > 1:
            p = paragraph._p
            p.remove(run._r)
    for run, target_text, color in zip(paragraph.runs, target_texts, colors):
        run.font.color.rgb = RGBColor(*color)
        run.text = target_text


def gen_barlength(data):
    if not "총 언급 횟수" in data.columns:
        return None
    standard = data.loc[1, "총 언급 횟수"]
    barlength = 18.56 / standard
    barlengths = [18.56]
    for i in range(2, 11):
        out = data.loc[i, "총 언급 횟수"] * barlength
        barlengths.append(out)
    return barlengths


def update_slide(slide_num, filename, *, year=None, month=None):
    year, month = database.get_year_month(year, month)
    data_dir = database.get_data_dir(year, month)
    data_path = os.path.join(data_dir, filename)
    data = pd.read_csv(data_path, index_col=0)
    barlengths = gen_barlength(data)

    textboxes = []
    pictures = []
    rectangles = []

    for shape in TEMPLATE.slides[slide_num].shapes:
        logging.LOGGER.info(
            f"{shape.name} leftx: {shape.left}, topy:{shape.top}, width:{shape.width}, height:{shape.height}"
        )
        if shape.name.find("직사각형") != -1:
            rectangles.append(shape)
        elif shape.name.find("TextBox") != -1:
            textboxes.append(shape)
        elif shape.name.find("그림") != -1 or shape.name.find(
                "Picture") != -1 or shape.name.find("Rank") != -1:
            pictures.append(shape)

    textboxes.sort(key=lambda x: x.top, reverse=False)
    pictures.sort(key=lambda x: x.top, reverse=False)
    rectangles.sort(key=lambda x: x.top, reverse=False)

    logging.LOGGER.info("text------")
    count_idx = 0
    color = (240, 240, 240)
    for shape in textboxes:
        whole_text = "".join(
            [r.text for r in shape.text_frame.paragraphs[0].runs])
        if slide_num == 0 and whole_text == "69":
            target_text = str(
                data.index[data["Vtuber"] == "키즈나 아이"].tolist()[0])
            logging.LOGGER.info(target_text)
            text_frame = shape.text_frame
            p = text_frame.paragraphs[0]
            replace_text_only(p, target_text, color)

        elif re.search(strings.MONTH_PATTERN, whole_text) is not None:
            target_text = whole_text.replace(str(TEMPLATE_MONTH), str(year))

            text_frame = shape.text_frame
            p = text_frame.paragraphs[0]
            replace_text_only(p, target_text, color)
        elif whole_text.find(strings.COUNT_PATTERN) != -1:
            if slide_num == 0:
                detail_count_text = data[data["Vtuber"] ==
                                         "키즈나 아이"]["단어(언급 수)"].values[0]
                logging.LOGGER.info(detail_count_text)
            else:
                count_idx += 1
                detail_count_text = data.loc[count_idx, "단어(언급 수)"]
            text_frame = shape.text_frame
            p = text_frame.paragraphs[0]
            if len(detail_count_text) > 44:
                detail_count_text = detail_count_text[:44] + "..."
            replace_text_only(p, detail_count_text, color)
        elif re.search(strings.VIDEO_COUNT_PATTERN, whole_text) is not None:
            video_id = data.index[count_idx]
            detail_count_text = str(data.loc[video_id, "Count"]) + " 회"
            count_idx += 1
            text_frame = shape.text_frame
            p = text_frame.paragraphs[0]
            replace_text_only(p, detail_count_text, color)

        logging.LOGGER.info(
            f"{shape.name} leftx: {shape.left}, topy:{shape.top}, width:{shape.width}, height:{shape.height}, text: {whole_text}"
        )

    count_idx = 0
    logging.LOGGER.info("pic------")
    for shape in pictures:
        logging.LOGGER.info(
            f"{shape.name} leftx: {shape.left}, topy:{shape.top}, width:{shape.width}, height:{shape.height}"
        )
        left = shape.left
        top = shape.top
        width = shape.width
        height = shape.height
        # Replace image:shape.width = barlengths[count_idx]
        shape_name = str(shape.name)
        if shape_name.find("Rank") != -1:
            rank_num = int(shape_name.split("Rank")[1])
            if slide_num == 4:
                video_id = data.index[rank_num - 1]
                pic_dir = f"{data_dir}/{paths.THUMBNAIL_DIRNAME}/{video_id}.png"
                replace_picture(TEMPLATE.slides[slide_num], shape, pic_dir)
            else:
                vtuber_name = data.loc[rank_num, "Vtuber"]
                pic_dir = f"{paths.PROFILEPIC_PATH}/{vtuber_name}.png"
                replace_picture(TEMPLATE.slides[slide_num], shape, pic_dir)
        elif slide_num == 0 and shape.name == "그림 8":
            replace_picture(TEMPLATE.slides[slide_num], shape,
                            f"{data_dir}/{paths.WORDCLOUD_SAVE_FILENAME}")

    count_idx = 0
    logging.LOGGER.info("shape------")
    for shape in rectangles:
        whole_text = "".join(
            [r.text for r in shape.text_frame.paragraphs[0].runs])
        if re.search(strings.INCREASE_PATTERN, whole_text) is not None:
            if slide_num == 0:
                count = data[data["Vtuber"] == "키즈나 아이"]["총 언급 횟수"].values[0]
                increase_count = data[data["Vtuber"] ==
                                      "키즈나 아이"]["전월 대비 순위"].values[0]
            else:
                count = data.loc[count_idx + 1, "총 언급 횟수"]
                increase_count = data.loc[count_idx + 1, "전월 대비 순위"]
                shape.width = Cm(barlengths[count_idx])
                count_idx += 1

            if increase_count.find("▼") != -1:
                increase_color = (0, 113, 197)
            elif increase_count.find("▲") != -1 or increase_count.find(
                    "NEW") != -1:
                increase_color = (192, 0, 0)
            elif increase_count.find("■") != -1:
                increase_color = (13, 13, 13)
            else:
                increase_color = (13, 13, 13)

            target_texts = [str(count) + " ", str(increase_count)]
            colors = [(13, 13, 13), increase_color]

            text_frame = shape.text_frame
            p = text_frame.paragraphs[0]
            replace_texts_with_diff_color(p, target_texts, colors)
        logging.LOGGER.info(
            f"{shape.name} leftx: {shape.left}, topy:{shape.top}, width:{shape.width}, height:{shape.height}, text: {whole_text}"
        )


def save_ppt(year=None, month=None):
    data_dir = database.get_data_dir(year, month)
    TEMPLATE.save(f"{data_dir}/{month}월 통계.pptx")
